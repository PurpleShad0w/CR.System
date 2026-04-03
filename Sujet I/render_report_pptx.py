#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""render_report_pptx.py (Templates Slides focus)

Objectif (nouvelle règle demandée) :
- Vous allez créer dans "Templates Slides.pptx" des variantes de slide avec 6,5,4,3,2,1,0 images.
- Le renderer doit choisir AUTOMATIQUEMENT la variante qui correspond exactement au nombre d'images disponibles.
- On conserve le placement habituel (remplacement des placeholders PICTURE), et on ne casse plus la sélection.

Principe d'implémentation :
1) On détecte les slides "variants" dans la template via slide_types_template_slides.json (tokens visibles).
2) Pour chaque slide variante, on calcule le nombre de "photo slots" = nombre de shapes PICTURE situées dans la zone photo (partie droite).
3) Au rendu d'une slide, si nb_images = N, on choisit la variante ayant slot_count == N si elle existe.
   - sinon fallback : plus petite variante >= N, sinon la plus grande.
4) On remplit les N slots PICTURE par remplacement d'image, crop-to-fill.
5) Si la variante choisie contient plus de slots que N (cas fallback), on supprime les slots non utilisés + compagnons superposés.

Note importante :
- On NE désactive PAS l'ombre sur les slots remplis.

CLI:
  python render_report_pptx.py --template <pptx> --assembled <json> --out <pptx> --slide-types <json>
"""

from __future__ import annotations

import argparse
import json
import re
import time
from copy import deepcopy
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Pt

from PIL import Image

ORANGE = RGBColor(255, 92, 25)
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

FONT_BODY = 'Aptos'
FONT_HEADER = 'Verdana'

PT_BODY = 15
PT_HEADER = 32
PT_LEGEND = 12

LEGEND_BOTTOM_MARGIN = 36000  # ~1mm


def load_json(path: Path) -> Any:
    return json.loads(path.read_text(encoding='utf-8', errors='replace'))


def normalize_whitespace(text: str) -> str:
    t = (text or '').replace('\r\n', '\n').replace('\r', '\n')
    t = re.sub(r"\n{3,}", "\n\n", t)
    return t.strip()


def strip_markdown(text: str) -> str:
    if not text:
        return ''
    t = text
    t = re.sub(r"\*\*(.*?)\*\*", r"\1", t)
    t = re.sub(r"\*(.*?)\*", r"\1", t)
    t = re.sub(r"^\s*#+\s*", "", t, flags=re.MULTILINE)
    t = t.replace('•', '-')
    return t.strip()


def clean_bullets(text: str, max_lines: int = 12) -> List[str]:
    if not text:
        return []
    lines = []
    for ln in normalize_whitespace(strip_markdown(text)).split('\n'):
        s = ln.strip()
        if not s:
            continue
        if s.startswith('- '):
            s = s[2:].strip()
        elif s.startswith('-'):
            s = s[1:].strip()
        if not s:
            continue
        lines.append(s)
    if len(lines) > max_lines:
        lines = lines[:max_lines-1] + ['…']
    return lines


def slide_text(slide) -> str:
    parts: List[str] = []
    for shape in slide.shapes:
        if getattr(shape, 'has_text_frame', False) and shape.has_text_frame:
            parts.append(shape.text_frame.text or '')
    return "\n".join(parts)


def detect_template_slide_variants(tpl: Presentation, slide_types_cfg: Dict[str, Any]) -> Dict[str, List[Any]]:
    """Return type -> list of slides (variants) matching that type."""
    types = (slide_types_cfg.get('types') or {})
    found: Dict[str, List[Any]] = {}
    for sl in tpl.slides:
        stxt = slide_text(sl)
        for tname, tcfg in types.items():
            toks = tcfg.get('detect_tokens') or []
            if toks and all(tok in stxt for tok in toks):
                found.setdefault(tname, []).append(sl)
    return found


def clone_slide(prs_out: Presentation, slide_in) -> Any:
    blank_layout = prs_out.slide_layouts[6]
    new_slide = prs_out.slides.add_slide(blank_layout)
    try:
        new_slide._element.get_or_add_bg()._set_bg(slide_in._element.get_or_add_bg())
    except Exception:
        pass
    for shape in slide_in.shapes:
        new_slide.shapes._spTree.insert_element_before(deepcopy(shape._element), 'p:extLst')
    return new_slide


def _remove_shape(slide, shape) -> None:
    try:
        slide.shapes._spTree.remove(shape._element)
    except Exception:
        pass


# ----------------------------
# Template artifacts
# ----------------------------

def remove_template_artifacts(slide) -> None:
    for sh in list(slide.shapes):
        if getattr(sh, 'has_text_frame', False) and sh.has_text_frame:
            txt = (sh.text_frame.text or '').strip()
            low = txt.lower()
            if (
                'info clé' in low or 'info cle' in low
                or 'texte détail' in low or 'texte detail' in low
                or 'user flow' in low
                or re.fullmatch(r"page\s*", txt, flags=re.IGNORECASE)
            ):
                _remove_shape(slide, sh)


def enforce_etat_des_lieux_gtb(slide) -> None:
    for sh in list(slide.shapes):
        if getattr(sh, 'has_text_frame', False) and sh.has_text_frame:
            if (sh.text_frame.text or '').strip() == 'Etat des lieux':
                sh.text_frame.text = 'Etat des lieux GTB'


# ----------------------------
# Geometry helpers (for unused slot cleanup in fallback cases)
# ----------------------------

def _rect(sh) -> Tuple[int, int, int, int]:
    return (int(getattr(sh, 'left', 0)), int(getattr(sh, 'top', 0)), int(getattr(sh, 'width', 0)), int(getattr(sh, 'height', 0)))


def _area(r: Tuple[int, int, int, int]) -> int:
    return max(0, r[2]) * max(0, r[3])


def _intersection(a: Tuple[int, int, int, int], b: Tuple[int, int, int, int]) -> int:
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ax2, ay2 = ax + aw, ay + ah
    bx2, by2 = bx + bw, by + bh
    ix1, iy1 = max(ax, bx), max(ay, by)
    ix2, iy2 = min(ax2, bx2), min(ay2, by2)
    if ix2 <= ix1 or iy2 <= iy1:
        return 0
    return (ix2 - ix1) * (iy2 - iy1)


def _remove_unused_slot_companions(slide, slot_shape) -> None:
    """Remove companions (gradients/frames) that overlap an UNUSED slot."""
    sr = _rect(slot_shape)
    sa = _area(sr)
    if sa <= 0:
        return
    to_remove: List[Any] = []
    for sh in list(slide.shapes):
        if sh is slot_shape:
            continue
        # never remove text boxes
        if getattr(sh, 'has_text_frame', False) and sh.has_text_frame:
            continue
        r = _rect(sh)
        a = _area(r)
        if a <= 0:
            continue
        # avoid nuking huge backgrounds
        if a > sa * 12.0:
            continue
        inter = _intersection(sr, r)
        if inter <= 0:
            continue
        # require overlap
        if (inter / sa) < 0.06 and (inter / float(min(sa, a))) < 0.18:
            continue
        to_remove.append(sh)
    for sh in to_remove:
        _remove_shape(slide, sh)


# ----------------------------
# Image path resolution + conversion
# ----------------------------

_img_cache: Dict[str, Optional[Path]] = {}
_conv_cache: Dict[Path, Path] = {}


def infer_repo_root(assembled_path: Path) -> Path:
    start = assembled_path.resolve()
    for cand in [start.parent] + list(start.parents):
        if (cand / 'process').exists() and (cand / 'input').exists():
            return cand
        if cand.name.lower() == 'process' and (cand.parent / 'input').exists():
            return cand.parent
    return assembled_path.parent.parent.parent


def resolve_image_path(img: str, base_dir: Path, repo_root: Path) -> Optional[Path]:
    if not img:
        return None
    if img in _img_cache:
        return _img_cache[img]
    p = Path(img)
    if p.is_absolute() and p.exists():
        _img_cache[img] = p
        return p
    cand = (base_dir / img).resolve()
    if cand.exists():
        _img_cache[img] = cand
        return cand
    bn = p.name
    cand2 = (base_dir / bn).resolve()
    if cand2.exists():
        _img_cache[img] = cand2
        return cand2
    for sub in ('process', 'input'):
        base = repo_root / sub
        if not base.exists():
            continue
        try:
            for hit in base.rglob(bn):
                if hit.is_file():
                    _img_cache[img] = hit
                    return hit
        except Exception:
            continue
    _img_cache[img] = None
    return None


def magic_type(p: Path) -> str:
    try:
        b = p.read_bytes()[:16]
    except Exception:
        return 'unknown'
    if len(b) >= 3 and b[0:3] == b'\xFF\xD8\xFF':
        return 'jpeg'
    if len(b) >= 8 and b[0:8] == b'\x89PNG\r\n\x1a\n':
        return 'png'
    return 'unknown'


def normalize_image_for_ppt(p: Path, tmp_dir: Path) -> Optional[Path]:
    if not p.exists():
        return None
    if p in _conv_cache:
        return _conv_cache[p]
    kind = magic_type(p)
    ext = p.suffix.lower()
    if kind == 'jpeg' and ext in ('.jpg', '.jpeg'):
        _conv_cache[p] = p
        return p
    if kind == 'png' and ext == '.png':
        _conv_cache[p] = p
        return p
    try:
        im = Image.open(p)
        if im.mode not in ('RGB', 'RGBA'):
            im = im.convert('RGB')
        out = tmp_dir / (p.stem + '.png')
        im.save(out, format='PNG', optimize=True)
        _conv_cache[p] = out
        return out
    except Exception:
        return None


def images_to_pairs(images: List[Any]) -> List[Tuple[str, str]]:
    out: List[Tuple[str, str]] = []
    for im in images or []:
        if isinstance(im, str):
            out.append((im, ''))
        elif isinstance(im, dict):
            out.append(((im.get('path') or ''), (im.get('caption') or '')))
    return [(p, c) for p, c in out if p]


# ----------------------------
# Slot helpers (authoritative = picture placeholders)
# ----------------------------

def slot_pictures(slide, slide_w: int, slide_h: int) -> List[Any]:
    """Return PICTURE placeholders that correspond to photo slots (right-side grid)."""
    pics: List[Any] = []
    for sh in slide.shapes:
        try:
            if sh.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue
        except Exception:
            continue
        l, t, w, h = _rect(sh)
        # right side only
        if l < int(slide_w * 0.45):
            continue
        # exclude tiny pictures (logos etc.)
        if w < int(slide_w * 0.08) or h < int(slide_h * 0.08):
            continue
        pics.append(sh)
    pics.sort(key=lambda sh: (int(getattr(sh, 'top', 0)), int(getattr(sh, 'left', 0))))
    return pics


def slot_count(slide, slide_w: int, slide_h: int) -> int:
    return len(slot_pictures(slide, slide_w, slide_h))


def choose_variant_exact(variants: List[Any], nb_images: int, slide_w: int, slide_h: int) -> Any:
    """Choose the variant whose slot_count matches nb_images exactly if possible.
    Fallback: smallest variant with >= nb_images, else largest.
    """
    annot = [(slot_count(sl, slide_w, slide_h), sl) for sl in variants]
    annot.sort(key=lambda x: x[0])
    # exact match
    for n, sl in annot:
        if n == nb_images:
            return sl
    # smallest >=
    for n, sl in annot:
        if n >= nb_images:
            return sl
    return annot[-1][1]


# ----------------------------
# Picture replace + crop
# ----------------------------

def replace_picture_image(slide, pic_shape, img_file: Path) -> bool:
    try:
        _image_part, rId = slide.part.get_or_add_image_part(str(img_file))
    except Exception:
        return False
    try:
        blip = pic_shape._element.blipFill.blip
        blip.set(qn('r:embed'), rId)
        return True
    except Exception:
        try:
            blip = pic_shape._element.xpath('.//a:blip')[0]
            blip.set(qn('r:embed'), rId)
            return True
        except Exception:
            return False


def crop_to_fill(pic_shape, img_path: Path, box_w: int, box_h: int) -> None:
    try:
        im = Image.open(img_path)
        w, h = im.size
        if not w or not h or not box_h:
            return
        img_ar = w / h
        box_ar = float(box_w) / float(box_h)
        pic_shape.crop_left = 0.0
        pic_shape.crop_right = 0.0
        pic_shape.crop_top = 0.0
        pic_shape.crop_bottom = 0.0
        if img_ar > box_ar:
            new_w = h * box_ar
            excess = max(0.0, w - new_w)
            frac = (excess / w) / 2.0 if w else 0.0
            pic_shape.crop_left = frac
            pic_shape.crop_right = frac
        else:
            new_h = w / box_ar
            excess = max(0.0, h - new_h)
            frac = (excess / h) / 2.0 if h else 0.0
            pic_shape.crop_top = frac
            pic_shape.crop_bottom = frac
    except Exception:
        return


# ----------------------------
# Logo extraction
# ----------------------------

def extract_build4use_logo(template: Presentation, tmp_dir: Path) -> Tuple[Optional[Path], Optional[Tuple[int, int, int, int]]]:
    sw = int(template.slide_width)
    sh = int(template.slide_height)
    candidates = []
    for sl in template.slides:
        for shp in sl.shapes:
            try:
                if shp.shape_type != MSO_SHAPE_TYPE.PICTURE:
                    continue
            except Exception:
                continue
            try:
                left, top, width, height = int(shp.left), int(shp.top), int(shp.width), int(shp.height)
            except Exception:
                continue
            if left <= int(sw * 0.25) and top >= int(sh * 0.80) and width <= int(sw * 0.25) and height <= int(sh * 0.15):
                try:
                    blob = shp.image.blob
                    candidates.append((left, top, width, height, blob))
                except Exception:
                    pass
    if not candidates:
        return None, None
    candidates.sort(key=lambda x: x[2] * x[3])
    left, top, width, height, blob = candidates[0]
    logo_path = tmp_dir / 'logo_build4use.png'
    try:
        logo_path.write_bytes(blob)
    except Exception:
        return None, None
    return logo_path, (left, top, width, height)


def ensure_logo(slide_out, logo_path: Optional[Path], geom: Optional[Tuple[int, int, int, int]]) -> None:
    if not logo_path or not geom or not logo_path.exists():
        return
    sw = int(slide_out.part.presentation.slide_width)
    sh = int(slide_out.part.presentation.slide_height)
    for shp in slide_out.shapes:
        try:
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                left, top = int(shp.left), int(shp.top)
                if left <= int(sw * 0.25) and top >= int(sh * 0.80):
                    return
        except Exception:
            continue
    left, top, width, height = geom
    try:
        slide_out.shapes.add_picture(str(logo_path), left, top, width=width, height=height)
    except Exception:
        pass


# ----------------------------
# Text styling
# ----------------------------

def set_run_style(run, *, font_name: str, font_size_pt: int, color: RGBColor) -> None:
    try:
        run.font.name = font_name
    except Exception:
        pass
    try:
        run.font.size = Pt(font_size_pt)
    except Exception:
        pass
    try:
        run.font.color.rgb = color
    except Exception:
        pass


def style_text_frame(tf, *, font_name: str, font_size_pt: int, color: RGBColor, align: Optional[int] = None) -> None:
    try:
        for p in tf.paragraphs:
            if align is not None:
                p.alignment = align
            try:
                p.font.name = font_name
                p.font.size = Pt(font_size_pt)
                p.font.color.rgb = color
            except Exception:
                pass
            for r in p.runs:
                set_run_style(r, font_name=font_name, font_size_pt=font_size_pt, color=color)
    except Exception:
        pass


def set_slide_title(slide_out, title: str) -> None:
    if not title:
        return
    for sh in slide_out.shapes:
        if getattr(sh, 'has_text_frame', False) and sh.has_text_frame:
            txt = sh.text_frame.text or ''
            if '{{SLIDE_TITLE}}' in txt or 'Titre section' in txt:
                sh.text_frame.clear()
                p = sh.text_frame.paragraphs[0]
                p.text = title
                p.alignment = PP_ALIGN.LEFT
                style_text_frame(sh.text_frame, font_name=FONT_HEADER, font_size_pt=PT_HEADER, color=ORANGE, align=PP_ALIGN.LEFT)
                return


def set_body_bullets(slide_out, body: str, *, max_lines: int) -> None:
    lines = clean_bullets(body, max_lines=max_lines)
    for sh in slide_out.shapes:
        if getattr(sh, 'has_text_frame', False) and sh.has_text_frame:
            txt = sh.text_frame.text or ''
            if '{{TEXTE_BULLETS}}' in txt or 'Texte Texte' in txt:
                tf = sh.text_frame
                for p in tf.paragraphs:
                    p.text = ''
                if not lines:
                    return
                for i, s in enumerate(lines):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = s
                    p.level = 0
                    p.alignment = PP_ALIGN.LEFT
                try:
                    tf.word_wrap = True
                    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                except Exception:
                    pass
                style_text_frame(tf, font_name=FONT_BODY, font_size_pt=PT_BODY, color=BLACK, align=PP_ALIGN.LEFT)
                return


# ----------------------------
# Legends (simple mapping)
# ----------------------------

def is_legend_placeholder(t: str) -> bool:
    if not t:
        return False
    nt = re.sub(r"\s+", " ", t.strip()).lower()
    return ('légende photo' in nt) or ('legende photo' in nt)


def legend_shapes(slide) -> List[Any]:
    leg: List[Any] = []
    for sh in slide.shapes:
        if getattr(sh, 'has_text_frame', False) and sh.has_text_frame:
            if is_legend_placeholder(sh.text_frame.text or ''):
                leg.append(sh)
    leg.sort(key=lambda sh: (int(getattr(sh, 'top', 0)), int(getattr(sh, 'left', 0))))
    return leg


def fill_legends(slide_out, used_slot_count: int, imgs: List[Any], *, fallback_title: str = '') -> None:
    pairs = images_to_pairs(imgs)
    legs = legend_shapes(slide_out)
    if not legs:
        return
    m = min(len(pairs), used_slot_count, len(legs))
    for i in range(m):
        cap = (pairs[i][1] or '').strip() or (fallback_title or '').strip()
        if not cap:
            continue
        if len(cap) > 90:
            cap = cap[:87].rstrip() + '…'
        sh = legs[i]
        tf = sh.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = cap
        p.alignment = PP_ALIGN.CENTER
        try:
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.BOTTOM
            tf.margin_bottom = LEGEND_BOTTOM_MARGIN
            tf.margin_top = 0
            tf.margin_left = 0
            tf.margin_right = 0
        except Exception:
            pass
        style_text_frame(tf, font_name=FONT_BODY, font_size_pt=PT_LEGEND, color=WHITE, align=PP_ALIGN.CENTER)
    # clear remaining
    for j in range(m, min(len(legs), 6)):
        try:
            legs[j].text_frame.text = ''
        except Exception:
            pass


# ----------------------------
# Fill images using exact variant
# ----------------------------

def fill_images(slide_out, imgs: List[Any], base_dir: Path, repo_root: Path, tmp_dir: Path, slide_w: int, slide_h: int) -> int:
    pairs = images_to_pairs(imgs)
    slots = slot_pictures(slide_out, slide_w, slide_h)
    k = min(len(pairs), len(slots))

    filled = 0
    for i in range(k):
        path, _cap = pairs[i]
        ip = resolve_image_path(path, base_dir, repo_root)
        if not ip:
            continue
        normp = normalize_image_for_ppt(ip, tmp_dir)
        if not normp:
            continue
        ok = replace_picture_image(slide_out, slots[i], normp)
        if not ok:
            continue
        crop_to_fill(slots[i], normp, int(getattr(slots[i], 'width', 0)), int(getattr(slots[i], 'height', 0)))
        filled += 1

    # fallback safety: if slots exist beyond k, remove them + companions
    for i in range(k, len(slots)):
        _remove_unused_slot_companions(slide_out, slots[i])
        _remove_shape(slide_out, slots[i])

    return filled


# ----------------------------
# Save
# ----------------------------

def safe_save(prs: Presentation, out_path: Path) -> Path:
    out_path.parent.mkdir(parents=True, exist_ok=True)
    try:
        prs.save(str(out_path))
        return out_path
    except PermissionError:
        stamp = time.strftime('%Y%m%d_%H%M%S')
        alt = out_path.with_name(out_path.stem + f'_{stamp}' + out_path.suffix)
        prs.save(str(alt))
        print(f"WARNING: target PPTX locked, wrote instead: {alt}")
        return alt


# ----------------------------
# Build deck
# ----------------------------

def build_deck(base_template: Path, assembled_path: Path, out_path: Path, slide_types_path: Path) -> None:
    base_tpl = Presentation(str(base_template))
    data = load_json(assembled_path)
    repo_root = infer_repo_root(assembled_path)
    base_dir = assembled_path.parent
    tmp_dir = base_dir / '_tmp_img'
    tmp_dir.mkdir(parents=True, exist_ok=True)

    cfg = load_json(slide_types_path) if slide_types_path.exists() else {'types': {}, 'defaults': {}}
    variants_by_type = detect_template_slide_variants(base_tpl, cfg)

    prs_out = Presentation()
    prs_out.slide_width = base_tpl.slide_width
    prs_out.slide_height = base_tpl.slide_height

    logo_path, logo_geom = extract_build4use_logo(base_tpl, tmp_dir)

    today = datetime.now().strftime('%d/%m/%Y')
    global_map = {
        '{{DATE}}': today,
        '{{VERSION}}': '1',
        '{{CONTACT_EMAIL}}': 'contact@build4use.eu',
    }

    sw, sh = int(prs_out.slide_width), int(prs_out.slide_height)

    def replace_placeholders(slide_out) -> None:
        for shp in slide_out.shapes:
            if getattr(shp, 'has_text_frame', False) and shp.has_text_frame:
                txt = shp.text_frame.text or ''
                repl = txt
                for k0, v0 in global_map.items():
                    if k0 in repl:
                        repl = repl.replace(k0, v0)
                if repl != txt:
                    shp.text_frame.text = repl

    def emit_content_slide(s: Dict[str, Any]):
        stype = (s.get('type') or 'CONTENT_TEXT').strip()
        title = (s.get('title') or '').strip()
        body = (s.get('bullets') or s.get('body') or '')
        imgs = s.get('images') or []
        nb_images = len(images_to_pairs(imgs))

        cand_list = variants_by_type.get(stype) or variants_by_type.get('CONTENT_TEXT_IMAGES') or variants_by_type.get('CONTENT_TEXT') or []
        if not cand_list:
            return

        # Choose exact variant by nb_images
        slide_in = choose_variant_exact(cand_list, nb_images, sw, sh)

        slide_out = clone_slide(prs_out, slide_in)
        replace_placeholders(slide_out)

        remove_template_artifacts(slide_out)
        enforce_etat_des_lieux_gtb(slide_out)
        ensure_logo(slide_out, logo_path, logo_geom)

        set_slide_title(slide_out, title)
        set_body_bullets(slide_out, body, max_lines=12 if nb_images else 14)

        filled_count = fill_images(slide_out, imgs, base_dir, repo_root, tmp_dir, sw, sh)
        fill_legends(slide_out, filled_count, imgs, fallback_title=title)

    slides_spec = data.get('slides')
    if isinstance(slides_spec, list):
        for s in slides_spec:
            if isinstance(s, dict) and s.get('type') != 'PART_DIVIDER':
                emit_content_slide(s)

    final_path = safe_save(prs_out, out_path)
    print('Wrote:', final_path)


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument('--template', required=True)
    ap.add_argument('--assembled', required=True)
    ap.add_argument('--out', required=True)
    ap.add_argument('--slide-types', required=True)
    args = ap.parse_args()
    build_deck(Path(args.template), Path(args.assembled), Path(args.out), Path(args.slide_types))


if __name__ == '__main__':
    main()
