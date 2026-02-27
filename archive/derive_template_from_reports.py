#!/usr/bin/env python3
# -*- coding: utf-8 -*-

import argparse
import hashlib
import json
import os
import re
import zipfile
from collections import Counter, defaultdict
from copy import deepcopy
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ---------- utilities ----------
def slide_signature(slide, tol=5000):
    """
    Build a coarse signature that tends to be stable across reports
    using same layout: (shape_type, rounded geometry, has_text)
    tol in EMU (approx tolerance).
    """
    sig = []
    for sh in slide.shapes:
        st = int(sh.shape_type)
        # geometry
        left = int(getattr(sh, "left", 0) // tol)
        top  = int(getattr(sh, "top", 0) // tol)
        w    = int(getattr(sh, "width", 0) // tol)
        h    = int(getattr(sh, "height", 0) // tol)
        has_text = int(getattr(sh, "has_text_frame", False) and sh.has_text_frame)
        sig.append((st, left, top, w, h, has_text))
    sig.sort()
    return hashlib.md5(repr(sig).encode("utf-8")).hexdigest()

def scrub_text_in_slide(slide):
    # Replace obvious case tokens; you can extend with more patterns
    patterns = [
        (r"\bP\\d{5,6}\\b", "{{CASE_CODE}}"),
        (r"\\bVersion\\s*[0-9.]+\\b", "Version {{VERSION}}"),
        (r"\\b(19|20)\\d{2}\\b", "{{YEAR}}"),
    ]
    for sh in slide.shapes:
        if getattr(sh, "has_text_frame", False) and sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for run in p.runs:
                    t = run.text
                    for pat, rep in patterns:
                        t = re.sub(pat, rep, t, flags=re.IGNORECASE)
                    run.text = t

def remove_picture_shapes(slide):
    for sh in list(slide.shapes):
        if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
            el = sh._element
            el.getparent().remove(el)

# ---------- background image stripping (the key fix) ----------
def strip_background_images(pptx_in: Path, pptx_out: Path):
    """
    Remove slide background images by editing ppt/slides/slideX.xml:
    delete <p:bg> nodes that contain <a:blipFill>.
    """
    with zipfile.ZipFile(pptx_in, "r") as zin:
        with zipfile.ZipFile(pptx_out, "w", compression=zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                if item.filename.startswith("ppt/slides/slide") and item.filename.endswith(".xml"):
                    # crude but effective: remove bg with blipFill
                    xml = data.decode("utf-8", errors="ignore")
                    # remove any <p:bg> ... </p:bg> block containing "blipFill"
                    xml = re.sub(r"<p:bg>.*?blipFill.*?</p:bg>", "", xml, flags=re.DOTALL)
                    data = xml.encode("utf-8")
                zout.writestr(item, data)

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--reports", nargs="+", required=True, help="List of source report PPTX paths")
    ap.add_argument("--out", required=True, help="Output template PPTX")
    ap.add_argument("--keep_top", type=int, default=6, help="How many most-common layouts to keep")
    args = ap.parse_args()

    # Collect signatures across all reports
    sig_to_examples = defaultdict(list)
    sig_counts = Counter()

    for rp in args.reports:
        prs = Presentation(rp)
        for i, slide in enumerate(prs.slides):
            sig = slide_signature(slide)
            sig_counts[sig] += 1
            if len(sig_to_examples[sig]) < 3:
                sig_to_examples[sig].append((rp, i))

    # pick most common slide signatures as baseline layouts
    common = [s for s, c in sig_counts.most_common(args.keep_top)]

    # Build template by copying representative slides (one each)
    base_report, base_idx = sig_to_examples[common[0]][0]
    base_prs = Presentation(base_report)

    template = Presentation()
    # match slide size to base
    template.slide_width = base_prs.slide_width
    template.slide_height = base_prs.slide_height

    for sig in common:
        src_path, slide_idx = sig_to_examples[sig][0]
        src_prs = Presentation(src_path)
        src_slide = src_prs.slides[slide_idx]

        # add blank slide with same layout index if possible (fallback to blank)
        try:
            layout = template.slide_layouts[0]
        except Exception:
            layout = template.slide_layouts[0]
        new_slide = template.slides.add_slide(layout)

        # copy shapes
        for sh in src_slide.shapes:
            new_slide.shapes._spTree.insert_element_before(deepcopy(sh._element), 'p:extLst')

        # scrub variable content
        remove_picture_shapes(new_slide)
        scrub_text_in_slide(new_slide)

    tmp = Path(args.out).with_suffix(".tmp.pptx")
    template.save(tmp)

    # crucial: remove background images too
    strip_background_images(tmp, Path(args.out))
    tmp.unlink(missing_ok=True)

    print(f"Template written: {args.out}")

if __name__ == "__main__":
    main()