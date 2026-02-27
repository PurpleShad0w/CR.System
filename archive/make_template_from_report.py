#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
Make a reusable PPTX template from an existing audit report PPTX.

Typical use:
  python make_template_from_report.py --in "P060011 ... v1.1.pptx" --out "TEMPLATE_AUDIT_BACS.pptx" --keep-parts 1 2 3 --strip-images

It will:
- keep cover + sommaire + slides belonging to parts 1-3 (based on the Sommaire slide ranges)
- remove other slides
- scrub text (client/site codes/dates) into placeholders
- optionally strip pictures
"""

import argparse
import re
from copy import deepcopy
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# ---------- slide delete helpers (python-pptx doesn't expose delete API) ----------
def delete_slide(prs, slide):
    slide_id = slide.slide_id
    slides = prs.slides._sldIdLst  # pylint: disable=protected-access
    for sldId in list(slides):
        if int(sldId.id) == slide_id:
            slides.remove(sldId)
            break

def delete_slides_by_indices(prs, indices_to_delete):
    # delete from end to start
    for idx in sorted(indices_to_delete, reverse=True):
        delete_slide(prs, prs.slides[idx])

# ---------- text scrubbing ----------
PLACEHOLDERS = [
    # generic placeholders
    (r"\bP\d{5,6}\b", "{{CASE_CODE}}"),
    (r"\bVersion\s*[0-9.]+\b", "Version {{VERSION}}"),
    (r"\b(19|20)\d{2}\b", "{{YEAR}}"),
    # months in French (very rough)
    (r"\b(janvier|février|mars|avril|mai|juin|juillet|août|septembre|octobre|novembre|décembre)\b", "{{MONTH}}"),
]

def scrub_text(s):
    if not s:
        return s
    out = s
    for pattern, repl in PLACEHOLDERS:
        out = re.sub(pattern, repl, out, flags=re.IGNORECASE)
    # common names you may want to replace
    # (add your client/site tokens here as you encounter them)
    # out = re.sub(r"Goussonville", "{{SITE}}", out, flags=re.IGNORECASE)
    return out

def scrub_slide(slide, strip_images=False):
    for shape in list(slide.shapes):
        # remove pictures if requested
        if strip_images and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            sp = shape._element
            sp.getparent().remove(sp)
            continue

        # scrub text frames
        if hasattr(shape, "text_frame") and shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.text = scrub_text(run.text)

# ---------- detect ranges from Sommaire ----------
RANGE_RE = re.compile(r"slides\\s+(\\d+)\\s+à\\s+(\\d+)", re.IGNORECASE)

def get_slide_text(slide):
    txt = []
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.has_text_frame:
            txt.append(shape.text_frame.text)
    return "\n".join(txt)

def infer_part_ranges_from_sommaire(prs):
    """
    Try to parse slide 2 (index 1) as Sommaire and extract ranges like 'slides 3 à 17'.
    Returns list of ranges in order of appearance.
    """
    if len(prs.slides) < 2:
        return []
    sommaire = prs.slides[1]
    text = get_slide_text(sommaire)
    ranges = []
    for m in RANGE_RE.finditer(text):
        a, b = int(m.group(1)), int(m.group(2))
        ranges.append((a, b))
    return ranges

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--in", dest="inp", required=True, help="Input report PPTX")
    ap.add_argument("--out", dest="out", required=True, help="Output template PPTX")
    ap.add_argument("--keep-parts", nargs="+", type=int, default=[1,2,3],
                    help="Macro parts to keep (1-based). Default: 1 2 3")
    ap.add_argument("--strip-images", action="store_true", help="Remove picture shapes")
    args = ap.parse_args()

    prs = Presentation(args.inp)

    # Sommaire often lists: part1 range, part2 range, part3 range, part4 range, ...
    ranges = infer_part_ranges_from_sommaire(prs)

    # Slide numbering in the Sommaire is 1-based "slides X à Y"
    # python-pptx indices are 0-based
    keep_slide_indices = set([0, 1])  # cover + sommaire

    # keep parts 1..3 by selecting ranges from sommaire
    for part_num in args.keep_parts:
        if part_num <= 0 or part_num > len(ranges):
            continue
        start, end = ranges[part_num - 1]
        # convert to indices
        for s in range(start-1, end):
            if 0 <= s < len(prs.slides):
                keep_slide_indices.add(s)

    # delete everything else
    delete_indices = [i for i in range(len(prs.slides)) if i not in keep_slide_indices]
    delete_slides_by_indices(prs, delete_indices)

    # scrub remaining slides
    for slide in prs.slides:
        scrub_slide(slide, strip_images=args.strip_images)

    prs.save(args.out)
    print(f"Template written: {args.out}")

if __name__ == "__main__":
    main()