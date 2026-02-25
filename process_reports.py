#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
process_reports.py

Purpose
-------
Ingest audit reports (PDF / DOCX / PPTX) and convert them into a normalized, LLM-ready intermediate representation.

Inputs
------
<root>/<audits>/**/*.(pdf|docx|pptx)

Outputs
-------
processed_reports/Rapports_d_audit/
  docs/<doc_id>.json
  chunks/<doc_id>.jsonl
  assets/<doc_id>/*
  manifest.json

Usage Example
--------------
python process_reports.py --out processed_reports --extract-images
"""

import argparse
import hashlib
import json
import re
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path


# ============================================================
# 1) Optional dependencies (fail gracefully)
# ============================================================

try:
    import PyPDF2
except Exception:
    PyPDF2 = None

try:
    import docx
except Exception:
    docx = None

try:
    from pptx import Presentation
except Exception:
    Presentation = None


# ============================================================
# 2) Regex & constants
# ============================================================

RE_MULTI_SPACE = re.compile(r"[ \t]{2,}")
RE_PAGE_HEADER_NOISE = re.compile(r"Page\s+\d+\s+sur\s+\d+", re.IGNORECASE)

# NOTE: These keywords remain intentionally French
PART1_KEYS = (
    "contexte", "présentation", "presentation", "généralités", "generalites",
    "objectifs", "périmètre", "perimetre", "visite de site"
)
PART2_KEYS = (
    "etat des lieux", "état des lieux", "architecture",
    "etat du système", "état du système",
    "etat du superviseur",
    "réseau", "reseau",
    "liste des équipements", "obsolescences"
)
PART3_KEYS = (
    "diagnostic", "diagnostics", "constats", "essais",
    "résultats", "resultats",
    "dysfonctionnements", "etat des communications"
)


# ============================================================
# 3) Core data model
# ============================================================

@dataclass
class Block:
    """
    Atomic extracted unit.

    type:
      - heading
      - slide_title
      - paragraph
      - table
      - notes
    """
    type: str
    text: str
    ref: str   # page or slide reference (e.g. p003, s012)


# ============================================================
# 4) Utilities
# ============================================================

def utc_now_iso() -> str:
    return datetime.now(timezone.utc).isoformat()

def doc_id_for_file(path: Path) -> str:
    """
    Stable document identifier:
    <sanitized filename>-<short sha1>
    """
    h = hashlib.sha1(str(path).encode("utf-8")).hexdigest()[:12]
    safe = re.sub(r"[^a-zA-Z0-9_-]+", "_", path.stem)[:60]
    return f"{safe}-{h}"

def clean_text(text: str) -> str:
    """
    Light text normalization.
    This is intentionally conservative.
    """
    t = text.replace("\u00a0", " ")
    t = RE_MULTI_SPACE.sub(" ", t)
    t = t.strip()
    t = RE_PAGE_HEADER_NOISE.sub("", t).strip()
    return t

def classify_part(title_or_text: str) -> int | None:
    """
    Rough heuristic to classify content into Part 1 / 2 / 3.
    Used later for chunk grouping and skeleton learning.
    """
    s = (title_or_text or "").lower()
    if any(k in s for k in PART1_KEYS):
        return 1
    if any(k in s for k in PART2_KEYS):
        return 2
    if any(k in s for k in PART3_KEYS):
        return 3
    return None


# ============================================================
# 5) PDF extraction
# ============================================================

def extract_pdf(path: Path) -> tuple[list[Block], dict]:
    """
    Extract text from PDF using PyPDF2.

    Strategy:
    - Uppercase lines and "Sommaire" are treated as headings
    - Everything else becomes paragraphs
    """
    if PyPDF2 is None:
        raise RuntimeError("PyPDF2 not installed")

    blocks: list[Block] = []
    meta = {"pages": None, "needs_ocr": False}

    with path.open("rb") as f:
        reader = PyPDF2.PdfReader(f)
        meta["pages"] = len(reader.pages)

        total_len = 0
        for i, page in enumerate(reader.pages, start=1):
            try:
                txt = page.extract_text() or ""
            except Exception:
                txt = ""

            txt = clean_text(txt)
            total_len += len(txt)

            for line in txt.splitlines():
                line = clean_text(line)
                if not line:
                    continue

                if line.lower().startswith("sommaire") or (len(line) > 6 and line.isupper()):
                    blocks.append(Block("heading", line, ref=f"p{i:03d}"))
                else:
                    blocks.append(Block("paragraph", line, ref=f"p{i:03d}"))

        if total_len < 500:
            meta["needs_ocr"] = True

    return blocks, meta


# ============================================================
# 6) DOCX extraction
# ============================================================

def extract_docx(path: Path) -> tuple[list[Block], dict]:
    """
    Extract text from Word documents using python-docx.
    """
    if docx is None:
        raise RuntimeError("python-docx not installed")

    d = docx.Document(str(path))
    blocks: list[Block] = []
    meta = {"tables": 0}

    for p in d.paragraphs:
        txt = clean_text(p.text)
        if not txt:
            continue

        style = (p.style.name or "").lower() if p.style else ""
        if "heading" in style or "titre" in style:
            blocks.append(Block("heading", txt, ref="docx"))
        else:
            blocks.append(Block("paragraph", txt, ref="docx"))

    for ti, table in enumerate(d.tables, start=1):
        meta["tables"] += 1
        rows = []
        for row in table.rows:
            cells = [clean_text(c.text) for c in row.cells]
            rows.append("\t".join(cells))
        table_text = "\n".join(r for r in rows if r.strip())
        if table_text:
            blocks.append(Block("table", table_text, ref=f"table{ti:02d}"))

    return blocks, meta


# ============================================================
# 7) PPTX extraction
# ============================================================

def extract_pptx(path: Path,
                 extract_images: bool,
                 assets_dir: Path) -> tuple[list[Block], dict, list[dict]]:
    """
    Extract content from PowerPoint presentations.

    Strategy:
    - Slide titles become 'slide_title'
    - All other text boxes become paragraphs
    - Notes slides are preserved
    """
    if Presentation is None:
        raise RuntimeError("python-pptx not installed")

    prs = Presentation(str(path))
    blocks: list[Block] = []
    assets: list[dict] = []
    meta = {"slides": len(prs.slides)}

    for si, slide in enumerate(prs.slides, start=1):
        ref = f"s{si:03d}"

        title = ""
        if slide.shapes.title and slide.shapes.title.text:
            title = clean_text(slide.shapes.title.text)
            if title:
                blocks.append(Block("slide_title", title, ref=ref))

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                txt = clean_text(shape.text)
                if txt and txt != title:
                    blocks.append(Block("paragraph", txt, ref=ref))

            # Extract images (optional, for multimodal later)
            if extract_images and getattr(shape, "shape_type", None) == 13:
                image = shape.image
                ext = image.ext
                img_name = f"{ref}_{len(assets)+1:03d}.{ext}"
                out_path = assets_dir / img_name
                out_path.write_bytes(image.blob)
                assets.append({
                    "type": "image",
                    "slide": si,
                    "file": str(out_path),
                    "name": img_name
                })

        if slide.has_notes_slide:
            notes = clean_text(slide.notes_slide.notes_text_frame.text or "")
            if notes:
                blocks.append(Block("notes", notes, ref=ref))

    return blocks, meta, assets


# ============================================================
# 8) Outline & chunk building
# ============================================================

def build_outline(blocks: list[Block]) -> list[dict]:
    """
    Build a flat outline from headings and slide titles.
    """
    outline = []
    for b in blocks:
        if b.type in ("heading", "slide_title"):
            outline.append({
                "title": b.text,
                "ref": b.ref,
                "suggested_part": classify_part(b.text)
            })
    return outline

def chunk_blocks(blocks: list[Block],
                 max_chars: int = 1200,
                 overlap: int = 150) -> list[dict]:
    """
    Chunk blocks into overlapping text segments.
    """
    chunks = []
    buf = ""
    refs = set()
    current_part = None

    def flush():
        nonlocal buf, refs, current_part
        txt = clean_text(buf)
        if txt:
            chunks.append({
                "text": txt,
                "refs": sorted(refs),
                "part": current_part
            })
        buf = ""
        refs = set()

    for b in blocks:
        if b.type in ("heading", "slide_title"):
            p = classify_part(b.text)
            if p is not None:
                current_part = p

        if not b.text:
            continue

        if len(buf) + len(b.text) + 2 > max_chars:
            if overlap > 0 and len(buf) > overlap:
                tail = buf[-overlap:]
                flush()
                buf = tail + "\n" + b.text + "\n"
            else:
                flush()
                buf = b.text + "\n"
        else:
            buf += b.text + "\n"

        refs.add(b.ref)

    flush()
    return chunks

def parts_first3(chunks: list[dict]) -> dict:
    """
    Extract Part 1 / Part 2 / Part 3 chunk groups.
    """
    parts = {1: [], 2: [], 3: []}
    for c in chunks:
        if c.get("part") in (1, 2, 3):
            parts[c["part"]].append({"refs": c["refs"], "text": c["text"]})
    return {
        "part_1": parts[1],
        "part_2": parts[2],
        "part_3": parts[3],
    }


# ============================================================
# 9) Main
# ============================================================

def main():
    ap = argparse.ArgumentParser(description="Process audit reports into normalized JSON.")
    ap.add_argument("--root", default=".", help="Project root")
    ap.add_argument("--audits", default="Rapports d'audit", help="Input folder")
    ap.add_argument("--out", default="processed_reports", help="Output folder")
    ap.add_argument("--extract-images", action="store_true", help="Extract PPTX images")
    args = ap.parse_args()

    root = Path(args.root).resolve()
    audits = (root / args.audits).resolve()
    out_root = (root / args.out / "Rapports_d_audit").resolve()

    docs_dir = out_root / "docs"
    chunks_dir = out_root / "chunks"
    assets_root = out_root / "assets"

    docs_dir.mkdir(parents=True, exist_ok=True)
    chunks_dir.mkdir(parents=True, exist_ok=True)
    assets_root.mkdir(parents=True, exist_ok=True)

    supported = {".pdf", ".docx", ".pptx"}
    files = [p for p in audits.rglob("*") if p.suffix.lower() in supported]

    manifest = {
        "collection": "Rapports_d_audit",
        "folder": str(audits),
        "generated_at": utc_now_iso(),
        "files": [],
        "errors": [],
    }

    if not audits.exists():
        raise SystemExit(f"ERROR: audits folder not found: {audits}")

    for path in files:
        doc_id = doc_id_for_file(path)
        assets_dir = assets_root / doc_id
        if args.extract_images:
            assets_dir.mkdir(parents=True, exist_ok=True)

        try:
            ext = path.suffix.lower()
            if ext == ".pdf":
                blocks, meta = extract_pdf(path)
                assets = []
            elif ext == ".docx":
                blocks, meta = extract_docx(path)
                assets = []
            elif ext == ".pptx":
                blocks, meta, assets = extract_pptx(path, args.extract_images, assets_dir)
            else:
                continue

            outline = build_outline(blocks)
            chunks = chunk_blocks(blocks)
            parts = parts_first3(chunks)

            doc_pack = {
                "doc_id": doc_id,
                "file": str(path),
                "file_type": ext.lstrip("."),
                "generated_at": utc_now_iso(),
                "meta": meta,
                "outline": outline,
                "parts_first3": parts,
                "assets": assets,
            }

            (docs_dir / f"{doc_id}.json").write_text(
                json.dumps(doc_pack, ensure_ascii=False, indent=2),
                encoding="utf-8"
            )

            with (chunks_dir / f"{doc_id}.jsonl").open("w", encoding="utf-8") as f:
                for i, c in enumerate(chunks, start=1):
                    f.write(json.dumps({
                        "doc_id": doc_id,
                        "chunk_id": f"c{i:05d}",
                        "text": c["text"],
                        "refs": c["refs"],
                        "part": c.get("part"),
                    }, ensure_ascii=False) + "\n")

            manifest["files"].append({
                "doc_id": doc_id,
                "file": str(path),
                "type": ext,
                "blocks": len(blocks),
                "chunks": len(chunks),
                "assets": len(assets),
                "needs_ocr": bool(meta.get("needs_ocr")) if isinstance(meta, dict) else False,
            })

        except Exception as e:
            manifest["errors"].append({"file": str(path), "error": str(e)})

    (out_root / "manifest.json").write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2),
        encoding="utf-8"
    )

    print(json.dumps({
        "collection": "Rapports_d_audit",
        "files_processed": len(manifest["files"]),
        "errors": len(manifest["errors"]),
        "manifest": str(out_root / "manifest.json"),
    }, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()