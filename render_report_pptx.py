#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
render_report_pptx.py
Render a final audit report PPTX from:
- TEMPLATE_AUDIT_BUILD4USE.pptx (6-slide role template)
- assembled_report.json (LLM output)

Usage:
  python render_report_pptx.py --template TEMPLATE_AUDIT_BUILD4USE.pptx --assembled assembled_report.json --out Rapport_Audit.pptx
"""

import argparse
import json
import os
import re
from copy import deepcopy
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, List, Tuple

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

# -----------------------------
# Helpers: load / normalize
# -----------------------------

def load_json(path: Path) -> Dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))

def strip_markdown(text: str) -> str:
    if not text:
        return ""
    t = text
    # remove bold/italic markers
    t = re.sub(r"\*\*(.*?)\*\*", r"\1", t)
    t = re.sub(r"\*(.*?)\*", r"\1", t)
    # remove headings like "###"
    t = re.sub(r"^\s*#+\s*", "", t, flags=re.MULTILINE)
    # normalize bullets
    t = t.replace("•", "-")
    return t.strip()

def normalize_whitespace(text: str) -> str:
    t = (text or "").replace("\r\n", "\n").replace("\r", "\n")
    t = re.sub(r"\n{3,}", "\n\n", t)
    return t.strip()

def split_text_for_slides(text: str, max_chars: int = 900, max_lines: int = 12) -> List[str]:
    """
    Split text into slide-sized chunks conservatively.
    Prefers splitting on blank lines, then sentence boundaries, then hard cut.
    """
    text = normalize_whitespace(strip_markdown(text))
    if not text:
        return [""]

    # Break into paragraphs
    paras = [p.strip() for p in text.split("\n\n") if p.strip()]
    chunks = []
    cur = ""

    def cur_lines(s: str) -> int:
        return len([ln for ln in s.split("\n") if ln.strip()])

    for p in paras:
        candidate = (cur + "\n\n" + p).strip() if cur else p
        if len(candidate) <= max_chars and cur_lines(candidate) <= max_lines:
            cur = candidate
            continue

        # flush current if not empty
        if cur:
            chunks.append(cur)
            cur = ""

        # paragraph itself may be too long -> split by sentences
        if len(p) > max_chars or cur_lines(p) > max_lines:
            sentences = re.split(r"(?<=[\.\!\?])\s+", p)
            s_cur = ""
            for s in sentences:
                s = s.strip()
                if not s:
                    continue
                cand2 = (s_cur + " " + s).strip() if s_cur else s
                if len(cand2) <= max_chars:
                    s_cur = cand2
                else:
                    if s_cur:
                        chunks.append(s_cur.strip())
                    s_cur = s
            if s_cur:
                chunks.append(s_cur.strip())
        else:
            cur = p

    if cur:
        chunks.append(cur)

    return chunks if chunks else [""]

# -----------------------------
# PPTX cloning utilities
# -----------------------------

def clone_slide(prs_out: Presentation, slide_in) -> Any:
    """
    Clone a slide by deep-copying its shape XML into a new blank slide.
    Keeps formatting from the template slide.
    """
    blank_layout = prs_out.slide_layouts[6]  # blank
    new_slide = prs_out.slides.add_slide(blank_layout)

    # copy background if present (template has solid fills; still copy element)
    try:
        new_slide._element.get_or_add_bg()._set_bg(slide_in._element.get_or_add_bg())
    except Exception:
        pass

    # copy shapes
    for shape in slide_in.shapes:
        new_slide.shapes._spTree.insert_element_before(deepcopy(shape._element), 'p:extLst')

    return new_slide

# -----------------------------
# Placeholder replacement
# -----------------------------

def replace_in_text_frame(tf, mapping: Dict[str, str]):
    # Replace in each run to preserve styling as much as possible
    for p in tf.paragraphs:
        for r in p.runs:
            for k, v in mapping.items():
                if k in r.text:
                    r.text = r.text.replace(k, v)

def replace_placeholders(slide, mapping: Dict[str, str]):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            replace_in_text_frame(shape.text_frame, mapping)

def find_shape_containing(slide, token: str):
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
            if token in shape.text_frame.text:
                return shape
    return None

def set_bullets_in_shape(shape, text: str, font_size_pt: int = 16):
    """
    Replace the text frame content with bullet paragraphs.
    Accepts text with '-' bullet lines or normal lines.
    """
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True

    text = normalize_whitespace(strip_markdown(text))
    lines = [ln.strip() for ln in text.split("\n") if ln.strip()]

    if not lines:
        p = tf.paragraphs[0]
        p.text = ""
        return

    for i, ln in enumerate(lines):
        is_bullet = ln.startswith("- ")
        content = ln[2:].strip() if is_bullet else ln

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = content
        p.level = 0
        p.font.size = Pt(font_size_pt)
        p.font.name = "Calibri"
        p.font.bold = False

        if is_bullet:
            p.space_before = Pt(2)
            p.space_after = Pt(2)
            p.alignment = PP_ALIGN.LEFT

# -----------------------------
# Build output deck
# -----------------------------

def build_deck(template_path: Path, assembled_path: Path, out_path: Path):
    tpl = Presentation(str(template_path))
    data = load_json(assembled_path)

    case_id = data.get("case_id") or assembled_path.stem
    report_type = data.get("report_type") or "AUDIT"
    today = datetime.now().strftime("%d/%m/%Y")

    # Macro parts in assembled_report.json are expected like:
    # {"macro_parts":[{"macro_part":1,"macro_part_name":"...","sections":[{"bucket_id":"...","text":"..."}]}]}
    macro_parts = data.get("macro_parts") or []

    # Determine titles for parts 1–3 (fallback to report_type defaults)
    part_titles = {1: "Partie 1", 2: "Partie 2", 3: "Partie 3"}
    for mp in macro_parts:
        try:
            mp_num = int(mp.get("macro_part"))
            mp_name = mp.get("macro_part_name") or f"Partie {mp_num}"
            part_titles[mp_num] = mp_name
        except Exception:
            pass

    # Output presentation, same size as template
    prs_out = Presentation()
    prs_out.slide_width = tpl.slide_width
    prs_out.slide_height = tpl.slide_height

    # Template slides by role (fixed in TEMPLATE_AUDIT_BUILD4USE.pptx)
    # 0 cover, 1 sommaire, 2 section header, 3 content+visual, 4 text-only, 5 conclusion
    s_cover_in = tpl.slides[0]
    s_toc_in = tpl.slides[1]
    s_section_in = tpl.slides[2]
    s_content_in = tpl.slides[3]
    s_text_in = tpl.slides[4]
    s_concl_in = tpl.slides[5]

    # ---------- Cover ----------
    s_cover = clone_slide(prs_out, s_cover_in)
    cover_map = {
        "{{AUDIT_TYPE}}": report_type,
        "{{CLIENT}}": "N/A",
        "{{SITE}}": case_id,
        "{{VILLE}}": "",
        "{{CASE_CODE}}": case_id,
        "{{DATE}}": today,
        "{{VERSION}}": "1",
        "{{CONTACT_EMAIL}}": "contact@build4use.eu"
    }
    replace_placeholders(s_cover, cover_map)

    # ---------- Sommaire ----------
    s_toc = clone_slide(prs_out, s_toc_in)
    toc_map = {
        "{{PART1_TITLE}}": part_titles.get(1, "Partie 1"),
        "{{PART2_TITLE}}": part_titles.get(2, "Partie 2"),
        "{{PART3_TITLE}}": part_titles.get(3, "Partie 3"),
        "{{PART4_TITLE}}": "",
        "{{PART5_TITLE}}": "",
        "{{DATE}}": today,
        "{{VERSION}}": "1",
        "{{CONTACT_EMAIL}}": "contact@build4use.eu"
    }
    replace_placeholders(s_toc, toc_map)

    # ---------- Content for macro-parts 1..3 ----------
    for mp_num in [1, 2, 3]:
        mp = next((x for x in macro_parts if int(x.get("macro_part", -1)) == mp_num), None)
        if not mp:
            continue

        mp_name = mp.get("macro_part_name") or f"Partie {mp_num}"

        # Section header slide
        s_sec = clone_slide(prs_out, s_section_in)
        sec_map = {
            "{{PART1_TITLE}}": mp_name,  # the template uses {{PART1_TITLE}} on this slide sample
            "{{DATE}}": today,
            "{{VERSION}}": "1",
            "{{CONTACT_EMAIL}}": "contact@build4use.eu"
        }
        replace_placeholders(s_sec, sec_map)

        # Also replace the number badge "01" manually if present
        # (search for a shape that contains '01' and replace with mp_num formatted)
        badge_num = f"{mp_num:02d}"
        for sh in s_sec.shapes:
            if getattr(sh, "has_text_frame", False) and sh.has_text_frame:
                if sh.text_frame.text.strip() in {"01", "02", "03", "99"}:
                    sh.text_frame.text = badge_num

        # Section slides
        for sec in mp.get("sections") or []:
            bucket_id = sec.get("bucket_id") or "SECTION"
            sec_text = sec.get("text") or ""
            chunks = split_text_for_slides(sec_text, max_chars=900, max_lines=12)

            for idx, chunk in enumerate(chunks, start=1):
                # choose layout: text-only by default (no images provided)
                slide_out = clone_slide(prs_out, s_text_in)

                # Title: bucket_id, with (1/2) if split
                title = bucket_id
                if len(chunks) > 1:
                    title = f"{bucket_id} ({idx}/{len(chunks)})"

                replace_placeholders(slide_out, {
                    "{{SLIDE_TITLE}}": title,
                    "{{DATE}}": today,
                    "{{VERSION}}": "1",
                    "{{CONTACT_EMAIL}}": "contact@build4use.eu"
                })

                # Fill the bullet placeholder
                shape = find_shape_containing(slide_out, "{{TEXTE_BULLETS}}")
                if shape:
                    set_bullets_in_shape(shape, chunk, font_size_pt=16)
                else:
                    # fallback: first textbox gets the content
                    for sh in slide_out.shapes:
                        if getattr(sh, "has_text_frame", False) and sh.has_text_frame:
                            set_bullets_in_shape(sh, chunk, font_size_pt=16)
                            break

    # ---------- Conclusion ----------
    s_concl = clone_slide(prs_out, s_concl_in)

    # naive conclusion fallback: pick last paragraphs from part 3 if available
    conclusion_text = ""
    mp3 = next((x for x in macro_parts if int(x.get("macro_part", -1)) == 3), None)
    if mp3 and (mp3.get("sections") or []):
        last_txt = mp3["sections"][-1].get("text") or ""
        last_txt = normalize_whitespace(strip_markdown(last_txt))
        conclusion_text = last_txt[:900]  # keep it short
    if not conclusion_text:
        conclusion_text = "Synthèse à compléter (données de conclusion non fournies dans l'export)."

    replace_placeholders(s_concl, {
        "{{CONCLUSION_TEXTE}}": conclusion_text,
        "{{DATE}}": today,
        "{{VERSION}}": "1",
        "{{CONTACT_EMAIL}}": "contact@build4use.eu"
    })

    prs_out.save(str(out_path))

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("--template", required=True, help="Path to TEMPLATE_AUDIT_BUILD4USE.pptx")
    ap.add_argument("--assembled", required=True, help="Path to assembled_report.json")
    ap.add_argument("--out", default="Rapport_Audit.pptx", help="Output PPTX path")
    args = ap.parse_args()

    build_deck(Path(args.template), Path(args.assembled), Path(args.out))
    print(f"Wrote: {args.out}")

if __name__ == "__main__":
    main()